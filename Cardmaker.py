import argparse, os
from datetime import date
from io import BytesIO

import qrcode
from PIL import Image
from pptx import Presentation
from pptx.util import *
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from mutagen.mp3 import MP3
from mutagen.id3 import ID3
from mutagen.easyid3 import EasyID3

from library.data import MusicLibrary
from library.system_setting import Setting


class CardMaker:
    CONTROLLER = [('play', 'A 1'),
    ('pause', 'A 2'),
    ('stop', 'A 3'),
    ('previous', 'A 4'),
    ('next', 'A 5'),
    ('playlist', 'A 6'),
    ('feeling lucky', 'A 7'),
    ('update', 'D 1'),]

    def __init__(self):
        self.default_wd = os.getcwd()
        self.ribbon_path = os.path.join(self.default_wd, 'ForgedCards/template/ribbon.png')
        self.current_year = date.today().year
        self.settings = Setting()
        os.chdir(self.settings.playlist_path)

    def generate_all(self):
        files = os.listdir(self.settings.playlist_path)

        for file in files:
            if file.endswith('.m3u'):
                print(file.replace('.m3u', ''))
                self.generate_deck(file.replace('.m3u', ''))

    def generate_deck(self, playlist_name):
        playlist_filepath = os.path.join(self.settings.playlist_path, playlist_name + '.m3u' )
        if not os.path.isfile(playlist_filepath):
            print('Playlist not found')
            return None
        #deal with logo
        with open(playlist_filepath, "r", encoding='utf-8') as f:
            content = f.readlines()
        music_filepaths = [x.strip() for x in content]

        logo_path = os.path.join(self.settings.playlist_path, playlist_name + '.png' )
        if not os.path.isfile(logo_path):
            logo_path = os.path.join(self.default_wd, 'ForgedCards/template/logo.png')

        prs = Presentation()
        prs.slide_width = Cm(6)
        prs.slide_height = Cm(9)
        blank_slide_layout = prs.slide_layouts[6]

        for controller, qr_text in self.CONTROLLER:

            front_controller_slide = prs.slides.add_slide(blank_slide_layout)
            self.add_album_cover_info(front_controller_slide, controller, '')
            front_controller_slide.shapes.add_picture(os.path.join(self.default_wd, 'ForgedCards/controller/'+ controller +'.png'), Cm(2.6), Cm(6.9), Cm(0.8), Cm(0.8))
            self.add_front_bottom_right(front_controller_slide, playlist_name, None, music_filepaths, logo_path)
            self.add_bottomleft(front_controller_slide, qr_text)

            qr_code = qrcode.make(qr_text)
            qr_code.save(os.path.join(self.default_wd, 'ForgedCards/tmp/qr_code.png'))
            back_controller_slide = prs.slides.add_slide(blank_slide_layout)
            self.add_QR_code(back_controller_slide)
            self.add_ribbon(back_controller_slide)
            self.add_SAMsJukebox(back_controller_slide)
            self.add_MusicInYourHand(back_controller_slide)

        # FRONT SLIDE
        front_slide_playlist = prs.slides.add_slide(blank_slide_layout)
        self.add_album_cover_info(front_slide_playlist, playlist_name, 'Playlist')
        # bottom left
        self.add_bottomleft(front_slide_playlist, '')
        # bottom right
        self.add_front_bottom_right(front_slide_playlist, playlist_name, 0, music_filepaths, logo_path)
        self.add_genre(front_slide_playlist, 'Playlist')

        # QR Code Side
        qr_slide = prs.slides.add_slide(blank_slide_layout)
        qr_code = qrcode.make(playlist_name + '.m3u')
        qr_code.save(os.path.join(self.default_wd, 'ForgedCards/tmp/qr_code.png'))
        self.add_QR_code(qr_slide)
        self.add_ribbon(qr_slide)
        self.add_SAMsJukebox(qr_slide)
        self.add_MusicInYourHand(qr_slide)

        n = 0
        for music_file in music_filepaths:
            if music_file.endswith(".mp3") and music_file.startswith('../'):
                music_file = music_file[3:]
                print (music_file)
                songpath = self.settings.library_path + '/' + music_file

                n = n+1
                song_info = MP3(songpath, ID3=EasyID3)
                albumartist = song_info['albumartist'][0]
                tracknumber = song_info['tracknumber'][0]
                song_date = song_info['date'][0]
                song_title = song_info['title'][0]
                song_genre = song_info['genre'][0]
                qr_code = qrcode.make(music_file)
                qr_code.save(os.path.join(self.default_wd, 'ForgedCards/tmp/qr_code.png'))

                tags = ID3(songpath)
                pict = tags.get("APIC:").data
                im = Image.open(BytesIO(pict))
                im.save(os.path.join(self.default_wd, 'ForgedCards/tmp/album_cover.png'))

                # FRONT SLIDE
                front_slide = prs.slides.add_slide(blank_slide_layout)
                self.add_album_cover_info(front_slide, albumartist, song_title)
                # bottom left
                self.add_bottomleft(front_slide, tracknumber + ' ' + song_date)
                # bottom right
                self.add_front_bottom_right(front_slide, playlist_name, n, music_filepaths, logo_path)
                self.add_genre(front_slide, song_genre)

                # QR Code Side
                qr_slide = prs.slides.add_slide(blank_slide_layout)
                self.add_QR_code(qr_slide)
                self.add_ribbon(qr_slide)
                self.add_SAMsJukebox(qr_slide)
                self.add_MusicInYourHand(qr_slide)

        prs.save(os.path.join(self.settings.playlist_path, playlist_name+ '.pptx'))

    def add_album_cover_info(self, front_slide, albumartist, song_title):
        front_slide.shapes.add_picture(os.path.normpath(os.path.join(self.default_wd, 'ForgedCards/tmp/album_cover.png')), Cm(0.4), Cm(0.4), Cm(5.2), Cm(5.2))

        artist_txBox = front_slide.shapes.add_textbox(Cm(0.4), Cm(5.8), Cm(5.2), Cm(1))
        artist_text_frame = artist_txBox.text_frame
        artist_text_frame.clear()
        artist_p = artist_text_frame.paragraphs[0]
        artist_run = artist_p.add_run()

        artist_run.text = albumartist.upper()
        artist_font = artist_run.font
        artist_font.name = 'Selawik'
        artist_font.size = Pt(8)
        artist_font.bold = True
        artist_font.italic = None
        artist_font.color.rgb = RGBColor(0, 0, 0)
        artist_text_frame.margin_top = Cm(0.1)
        artist_text_frame.margin_bottom = Cm(0.1)
        artist_text_frame.margin_right = Cm(0)
        artist_text_frame.margin_left = Cm(0)

        song_txBox = front_slide.shapes.add_textbox(Cm(0.4), Cm(6.3), Cm(5.2), Cm(1))
        song_text_frame = song_txBox.text_frame
        song_text_frame.clear()
        song_p = song_text_frame.paragraphs[0]
        song_run = song_p.add_run()

        song_run.text = ' '.join(elem.capitalize() for elem in song_title.split())
        song_font = song_run.font
        song_font.name = 'Selawik'
        song_font.size = Pt(7)
        song_font.bold = False
        song_font.italic = None
        song_font.color.rgb = RGBColor(0, 0, 0)
        song_text_frame.margin_top = Cm(0.1)
        song_text_frame.margin_bottom = Cm(0.1)
        song_text_frame.margin_right = Cm(0)
        song_text_frame.margin_left = Cm(0)

    def add_bottomleft(self, front_slide, text):
        bottom_left_txBox = front_slide.shapes.add_textbox(Cm(0.4), Cm(8.4), Cm(2.6), Cm(0.5))
        bottom_left_text_frame = bottom_left_txBox.text_frame
        bottom_left_text_frame.clear()
        bottom_left_p = bottom_left_text_frame.paragraphs[0]
        bottom_left_run = bottom_left_p.add_run()

        bottom_left_run.text = text + ' SAM\'S JUKEBOX'
        bottom_left_font = bottom_left_run.font
        bottom_left_font.name = 'Selawik'
        bottom_left_font.size = Pt(5)
        bottom_left_font.color.rgb = RGBColor(0, 0, 0)
        bottom_left_text_frame.margin_top = Cm(0.1)
        bottom_left_text_frame.margin_bottom = Cm(0.1)
        bottom_left_text_frame.margin_right = Cm(0)
        bottom_left_text_frame.margin_left = Cm(0)

    def add_front_bottom_right(self, front_slide, playlist_name, n, music_filepaths, logo_path):
        # bottom right
        bottom_right_txBox = front_slide.shapes.add_textbox(Cm(2.9), Cm(8.4), Cm(2.6), Cm(0.5))
        bottom_right_text_frame = bottom_right_txBox.text_frame
        bottom_right_text_frame.clear()
        bottom_right_p = bottom_right_text_frame.paragraphs[0]
        bottom_right_run = bottom_right_p.add_run()

        text_b = playlist_name  + ' – ' + str(self.current_year)
        if n is not None:
            text_b = text_b + ' ' + str('%02d'%(n)) + '/' + str('%02d'%(len(music_filepaths)))
        bottom_right_run.text = text_b
        bottom_right_font = bottom_right_run.font
        bottom_right_font.name = 'Selawik'
        bottom_right_font.size = Pt(5)
        bottom_right_font.color.rgb = RGBColor(0, 0, 0)
        bottom_right_text_frame.margin_top = Cm(0.1)
        bottom_right_text_frame.margin_bottom = Cm(0.1)
        bottom_right_text_frame.margin_right = Cm(0)
        bottom_right_text_frame.margin_left = Cm(0)
        bottom_right_p.alignment = PP_ALIGN.RIGHT

        # Set logo
        front_slide.shapes.add_picture(logo_path, Cm(4.4), Cm(7.4), Cm(1), Cm(1))

    def add_QR_code(self, slide):
        slide.shapes.add_picture(os.path.join(self.default_wd, 'ForgedCards/tmp/qr_code.png'), Cm(0), Cm(0), Cm(6), Cm(6))

    def add_ribbon(self, slide):
        slide.shapes.add_picture(self.ribbon_path, Cm(0), Cm(7.6))

    def add_MusicInYourHand(self, slide):
        txBox = slide.shapes.add_textbox(Cm(0.3), Cm(7.1), Cm(5.40), Cm(0.5))
        text_frame = txBox.text_frame
        text_frame.clear()  # not necessary for newly-created shape
        p = text_frame.paragraphs[0]
        run = p.add_run()
        p.margin_top = Cm(0.13)
        p.margin_bottom = Cm(0.13)
        p.margin_right = Cm(0.25)
        p.margin_left = Cm(0.25)

        run.text = 'T H E   M U S I C   I N   Y O U R   H A N D S'
        font = run.font
        font.name = 'Roboto Condensed'
        font.size = Pt(6)
        font.bold = False
        font.italic = None
        font.color.rgb = RGBColor(66, 66, 78)

        p.alignment = PP_ALIGN.CENTER

    def add_SAMsJukebox(self, slide):
        txBox = slide.shapes.add_textbox(Cm(0.3), Cm(6.05), Cm(5.40), Cm(1.05))
        text_frame = txBox.text_frame
        text_frame.clear()  # not necessary for newly-created shape
        p = text_frame.paragraphs[0]
        run = p.add_run()
        p.margin_top = Cm(0.13)
        p.margin_bottom = Cm(0.13)
        p.margin_right = Cm(0.25)
        p.margin_left = Cm(0.25)
        run.text = 'SAM\' JUKEBOX'
        font = run.font
        font.name = 'Yellowtail'
        font.size = Pt(18)
        font.bold = False
        font.italic = None
        font.color.rgb = RGBColor(66, 66, 78)
        p.alignment = PP_ALIGN.CENTER

    def add_genre(self, front_slide, song_genre):
        # Set logo
        img_path = ''
        if song_genre in ['Pop', 'Dance-Pop']:
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Pop.png')
        elif song_genre == 'R&B':
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Rnb.png')
        elif song_genre == 'Rap/Hip Hop':
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Rap.png')
        elif song_genre == 'Contemporary Christian':
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Contemporary-christian.png')
        elif song_genre in ['Rock', 'Metal']:
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Rock.png')
        elif song_genre == 'Alternative':
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Alternative.png')
        elif song_genre == 'Dance':
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Dance.png')
        elif song_genre in ['Electro', 'Electronic', 'Electronica']:
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Electro.png')
        elif song_genre in ['Jazz']:
            img_path = os.path.join(self.default_wd, 'ForgedCards/template/musicgenre/Jazz.png')
        elif song_genre == 'Playlist':
            img_path = os.path.join(self.default_wd, 'ForgedCards/controller/Playlist.png')
        else:
            print (song_genre)
        if img_path != '':
            pic = front_slide.shapes.add_picture(img_path, Cm(0.4), Cm(7.4), Cm(1), Cm(1))

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Welcome to Sam\'s Jukebox Cards Maker")
    parser.add_argument('-S', '--Playlist_Name', type=str, default=None, help='Please provide the playlist name')
    args = parser.parse_args()

    cardmaker = CardMaker()
    if args.Playlist_Name:
        cardmaker.generate_deck(args.Playlist_Name)
    else:
        cardmaker.generate_all()
